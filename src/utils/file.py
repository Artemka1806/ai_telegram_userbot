import os
import tempfile
from pathlib import Path
import subprocess
import shutil
from src.utils.logger import logger
from src.config import Config
import mimetypes

# Make sure temp directories exist
os.makedirs(Config.TEMP_DIR, exist_ok=True)

async def convert_to_pdf(input_path, output_path=None):
    """Convert various file types to PDF format"""
    try:
        # Get file extension and mime type
        file_ext = os.path.splitext(input_path)[1].lower()
        mime_type = mimetypes.guess_type(input_path)[0]
        
        # If no output path specified, create one in temp directory
        if not output_path:
            output_path = os.path.join(Config.TEMP_DIR, f"{os.path.basename(input_path)}.pdf")
        
        logger.info(f"Converting file: {input_path} (type: {mime_type}) to PDF: {output_path}")
        
        # Handle different file types
        if file_ext in ['.pdf']:
            # Already PDF, just copy
            shutil.copy2(input_path, output_path)
            logger.info(f"File already in PDF format, copied to {output_path}")
            return output_path
            
        elif file_ext in ['.docx', '.doc']:
            # For Word documents, use docx2pdf
            from docx2pdf import convert
            convert(input_path, output_path)
            logger.info(f"Converted Word document to PDF: {output_path}")
            return output_path
            
        elif file_ext in ['.xlsx', '.xls']:
            # For Excel, use pandas and openpyxl
            import pandas as pd
            from fpdf import FPDF
            
            # Create PDF
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Arial", size=12)
            
            # Read all sheets
            xl = pd.ExcelFile(input_path)
            for sheet_name in xl.sheet_names:
                df = pd.read_excel(input_path, sheet_name=sheet_name)
                
                # Add sheet name as header
                pdf.cell(200, 10, txt=f"Sheet: {sheet_name}", ln=True, align='L')
                pdf.ln(5)
                
                # Add column headers
                for col in df.columns:
                    pdf.cell(40, 10, txt=str(col)[:20], border=1)
                pdf.ln()
                
                # Add rows (limit to first 1000 rows to prevent huge PDFs)
                for _, row in df.head(1000).iterrows():
                    for item in row:
                        pdf.cell(40, 10, txt=str(item)[:20], border=1)
                    pdf.ln()
                
                pdf.add_page()
            
            # Save PDF
            pdf.output(output_path)
            logger.info(f"Converted Excel to PDF: {output_path}")
            return output_path
            
        elif file_ext in ['.txt']:
            # For text files, use FPDF directly
            from fpdf import FPDF
            
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Arial", size=12)
            
            # Open text file and add content to PDF
            with open(input_path, 'r', encoding='utf-8', errors='ignore') as file:
                for line in file:
                    pdf.cell(0, 10, txt=line.strip(), ln=True)
            
            # Save PDF
            pdf.output(output_path)
            logger.info(f"Converted text file to PDF: {output_path}")
            return output_path
            
        elif file_ext in ['.pptx', '.ppt']:
            # For PowerPoint, use python-pptx and FPDF
            from pptx import Presentation
            from fpdf import FPDF
            
            prs = Presentation(input_path)
            pdf = FPDF()
            
            for slide in prs.slides:
                pdf.add_page()
                pdf.set_font("Arial", size=16)
                
                # Add slide title
                if slide.shapes.title:
                    pdf.cell(0, 10, txt=slide.shapes.title.text, ln=True)
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        pdf.set_font("Arial", size=12)
                        pdf.multi_cell(0, 10, txt=shape.text)
            
            # Save PDF
            pdf.output(output_path)
            logger.info(f"Converted PowerPoint to PDF: {output_path}")
            return output_path
            
        else:
            logger.warning(f"Unsupported file format for conversion: {file_ext}")
            return None
            
    except Exception as e:
        logger.error(f"Error converting file to PDF: {str(e)}")
        logger.exception(e)
        return None

async def process_file(file_path):
    """Process a file for analysis: convert to PDF if necessary"""
    try:
        file_ext = os.path.splitext(file_path)[1].lower()
        
        if file_ext == '.pdf':
            # Already a PDF file, return as is
            return file_path
        else:
            # Convert to PDF
            pdf_path = await convert_to_pdf(file_path)
            return pdf_path
    except Exception as e:
        logger.error(f"Error processing file: {str(e)}")
        logger.exception(e)
        return None